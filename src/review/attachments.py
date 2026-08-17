"""Attachment-directory indexing, evidence extraction, and reference matching."""
import logging
import re
import shutil
import subprocess
import zipfile
from collections import defaultdict
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, Iterable, List, Optional, Sequence, Tuple
from xml.etree import ElementTree

import openpyxl
from chardet import detect

from review.constants import CHECKPOINT_VOCAB
from review.embedded_media import (
    extract_docx_media,
    extract_pdf_media,
    extract_pptx_media,
)

_logger = logging.getLogger("review.attachments")
from review.excel_utils import (
    _build_sheet_text_for_llm,
    _extract_sheet_text_cells,
    _normalize_sheet_id,
    _truncate,
)
from review.finding_taxonomy import deterministic_finding_fields
from review.models import AttachmentFile, AttachmentPreviewItem, Finding

ATTACHMENT_FILE_RE = re.compile(
    r"([0-9A-Za-z_\-\.一-鿿]+?\.(?:png|jpg|jpeg|jp2|webp|gif|bmp|pdf|xlsx|xls|docx|doc|pptx|ppt|txt|csv|json|xml|log|md|html|htm|zip|tar|7z|rar|eml|msg))",
    re.IGNORECASE,
)
ATTACHMENT_PATH_RE = re.compile(
    r"([0-9A-Za-z_\-\.一-鿿]+(?:[\\/][0-9A-Za-z_\-\.一-鿿]+)+\.(?:png|jpg|jpeg|jp2|webp|gif|bmp|pdf|xlsx|xls|docx|doc|pptx|ppt|txt|csv|json|xml|log|md|html|htm|zip|tar|7z|rar|eml|msg))",
    re.IGNORECASE,
)
ATTACHMENT_INDEX_RE = re.compile(r"(?:附件|证据|图片|截图|索引|目录索引)\s*([0-9]{1,3})")

# Cross-sheet reference: matches `<C22.SA-3-1>`, `（C22.SA-3-1）`, `[C22.SA-3-1]`,
# `C22.SA-3-1` etc. The captured group is the dotted identifier; the leading
# `C22.` (workbook prefix) is optional in practice and is stripped. The
# remainder (`SA-3-1` / `PE-6` / `NS-5.3` …) is normalised to a directory
# key and matched against the attachment index's `by_sheet_norm` map so we
# can pull cross-sheet evidence (e.g. PE-6's `<C22.SA-3-1> 备份制度` from the
# SA-3-1 directory) into the current sheet's LLM context.
CROSS_SHEET_REF_RE = re.compile(
    r"[<\[]+\s*(?:C\d{1,3}\s*\.\s*)?((?:SA|PM|PE|NS)\s*[-.]?\s*\d{1,2}(?:[-.\s]\d{1,3})?[A-Za-z]?)\s*[>\]]+",
    re.IGNORECASE,
)
# Bare-prefix form: `C22.PE-6-2`, `C22.PE-6`, `C22.SA-3-1` (no surrounding
# brackets). This is the common workpaper convention and was missing entirely
# from the previous single-pattern regex.
CROSS_SHEET_BARE_RE = re.compile(
    r"\bC\d{1,3}\s*\.\s*((?:SA|PM|PE|NS)\s*[-.]?\s*\d{1,2}(?:[-.\s]\d{1,3})?[A-Za-z]?)\b",
    re.IGNORECASE,
)

_SHEET_TAG_RE = re.compile(r"\b((?:SA|PM|PE|NS)[-_ ]?\d{1,2}(?:[-_ ]\d{1,2})?[A-Za-z]?)\b", re.IGNORECASE)
_SHEET_TAG_NODELIM_RE = re.compile(r"\b((?:SA|PM|PE|NS)\d{1,2}(?:\d{1,2})?[A-Za-z]?)\b", re.IGNORECASE)

_TEXT_EXTENSIONS = {".txt", ".csv", ".json", ".xml", ".log", ".md", ".html", ".htm"}
_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".jp2", ".webp", ".gif", ".bmp"}
_MAX_ATTACHMENT_FILES = 500
_MAX_EXTRACTED_TEXT = 12000
_IGNORED_ATTACHMENT_NAMES = {".DS_Store", "Thumbs.db"}


def _normalize_rel_path(value: str) -> str:
    raw = str(value or "").strip().replace("\\", "/")
    if not raw or raw.startswith("/"):
        return ""
    parts = [part for part in raw.split("/") if part not in {"", "."}]
    if any(part == ".." for part in parts):
        return ""
    return "/".join(parts).lower()


def _extract_indices(text: str) -> List[str]:
    values: List[str] = []
    for value in ATTACHMENT_INDEX_RE.findall(str(text or "")):
        normalized = str(value).strip()
        if normalized and normalized not in values:
            values.append(normalized)
    return values


def _limit_text(text: str, limit: int = _MAX_EXTRACTED_TEXT) -> str:
    normalized = str(text or "").replace("\x00", "").replace("\r\n", "\n").replace("\r", "\n").strip()
    return normalized[:limit]


def _read_text_file(path: Path) -> Tuple[str, str]:
    raw = path.read_bytes()
    encoding = "utf-8"
    try:
        guessed = detect(raw).get("encoding")
        if guessed:
            encoding = str(guessed)
    except Exception:
        pass
    return _limit_text(raw.decode(encoding, errors="replace")), "ok"


def _read_xlsx_file(path: Path) -> Tuple[str, str]:
    workbook = openpyxl.load_workbook(path, data_only=True, read_only=True)
    lines: List[str] = []
    for worksheet in workbook.worksheets:
        for row in worksheet.iter_rows():
            for cell in row:
                if cell.value is None or str(cell.value).strip() == "":
                    continue
                lines.append(f"{worksheet.title}!{cell.coordinate}: {cell.value}")
                if sum(len(line) + 1 for line in lines) >= _MAX_EXTRACTED_TEXT:
                    return _limit_text("\n".join(lines)), "ok"
    return _limit_text("\n".join(lines)), "ok"


def _read_docx_file(path: Path) -> Tuple[str, str]:
    with zipfile.ZipFile(path) as archive:
        xml = archive.read("word/document.xml")
    root = ElementTree.fromstring(xml)
    text = "\n".join(
        value.strip()
        for node in root.iter()
        for value in [node.text or ""]
        if value.strip()
    )
    return _limit_text(text), "ok"


def _read_pptx_file(path: Path) -> Tuple[str, str]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    lines: List[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            value = getattr(shape, "text", "")
            if value and value.strip():
                lines.append(f"Slide {slide_number}: {value.strip()}")
    return _limit_text("\n".join(lines)), "ok"


def _read_pdf_file(path: Path) -> Tuple[str, str]:
    # pypdf is optional so the core project can still index directories when
    # the deployment image has only the standard document dependencies.
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        text = _limit_text("\n".join(page.extract_text() or "" for page in reader.pages))
        return text, "ok" if text else "binary"
    except ImportError:
        pdftotext = shutil.which("pdftotext")
        if not pdftotext:
            return "", "unavailable"
        result = subprocess.run(
            [pdftotext, "-layout", str(path), "-"],
            capture_output=True,
            text=True,
            timeout=30,
            check=False,
        )
        if result.returncode != 0:
            return "", "unavailable"
        text = _limit_text(result.stdout)
        return text, "ok" if text else "binary"
    except Exception:
        return "", "unavailable"


def _extract_attachment_text(path: Path) -> Tuple[str, str]:
    suffix = path.suffix.lower()
    try:
        if suffix in _TEXT_EXTENSIONS:
            return _read_text_file(path)
        if suffix in {".xls", ".doc"}:
            from review.legacy_convert import convert_legacy_to_modern
            converted = convert_legacy_to_modern(path)
            if converted is None:
                return "", "unsupported"
            modern_suffix = converted.suffix.lower()
            if modern_suffix == ".xlsx":
                return _read_xlsx_file(converted)
            if modern_suffix == ".docx":
                return _read_docx_file(converted)
            return "", "unsupported"
        if suffix == ".xlsx":
            return _read_xlsx_file(path)
        if suffix == ".docx":
            return _read_docx_file(path)
        if suffix == ".pptx":
            return _read_pptx_file(path)
        if suffix == ".pdf":
            return _read_pdf_file(path)
        if suffix in _IMAGE_EXTENSIONS:
            return "", "binary"
    except Exception:
        return "", "unavailable"
    return "", "unsupported"


@dataclass
class EvidenceEntry:
    rel_path: str
    file_type: str
    status: str
    excerpt: str
    source_document: Optional[str] = None
    is_embedded: bool = False


EVIDENCE_GUIDANCE = (
    "重要：[证据清单] 段列出本sheet附件目录中真实可用的文件及其嵌入图。\n"
    "若执行描述引用了「《某文档》」，**必须**从 [证据清单] 找出对应路径写入 evidence_refs.attachment 字段——不要只用文档名引用。\n"
    "若证据是截图（密码策略截图、系统参数界面），DOCX/PPTX 中抽取的嵌入图位于 .embedded_media/ 路径，格式 .embedded_media/<原文档>::<图名>.<ext>。\n"
    "不要把 [证据清单] 中不存在的文件写进 evidence_refs。\n"
    "不要因为「执行描述里没明说截图」就判证据不足——先看 [证据清单] 中是否真的缺。\n"
    "evidence_refs.attachment 字段是必填结构化字段，每条 evidence_ref 必须填实际文件路径或 .embedded_media/ 路径。\n"
)


def build_evidence_inventory(
    attachments: Optional[Dict[str, object]],
    *,
    max_entries: int = 30,
    max_embedded: int = 12,
    max_excerpt_chars: int = 200,
) -> str:
    """Build a compact evidence inventory for V1 review prompts.

    Lists real attachments (with status + excerpt) and embedded media
    (grouped by source document). Capped to keep prompt size bounded.
    Returns "" if attachments is empty/None.
    """
    if not attachments:
        return ""
    items = attachments.get("items", []) or []
    if not items:
        return ""

    real = []
    embedded_by_source: Dict[str, List[str]] = {}
    for it in items:
        rel = getattr(it, "rel_path", None) if not isinstance(it, dict) else it.get("rel_path")
        status = getattr(it, "status", "") if not isinstance(it, dict) else it.get("status", "")
        file_type = getattr(it, "file_type", "") if not isinstance(it, dict) else it.get("file_type", "")
        excerpt = getattr(it, "extracted_text", "") if not isinstance(it, dict) else it.get("extracted_text", "")
        if not rel:
            continue
        if rel.startswith(".embedded_media/") and "::" in rel:
            after = rel[len(".embedded_media/"):]
            source, media_name = after.split("::", 1)
            embedded_by_source.setdefault(source, []).append(media_name)
        else:
            excerpt_short = (excerpt or "")[:max_excerpt_chars].replace("\n", " ").strip()
            real.append((rel, status or "unknown", excerpt_short))

    total_real = len(real)
    total_embedded = sum(len(v) for v in embedded_by_source.values())
    real = real[:max_entries]
    embedded_pairs = []
    for src, names in embedded_by_source.items():
        for n in names:
            embedded_pairs.append((src, n))
    embedded_pairs = embedded_pairs[:max_embedded]

    lines = [f"[证据清单（前 {min(total_real, max_entries)} 个附件 + 前 {min(total_embedded, max_embedded)} 张嵌入图，目录实际有 {total_real} 项 + {total_embedded} 张）]\n"]
    lines.append(f"== 真实附件（{total_real} 项，列出前 {len(real)}） ==")
    for rel, status, ex in real:
        ex_part = f" — 摘录: {ex}" if ex else ""
        lines.append(f"[{status}] {rel}{ex_part}")
    lines.append("")
    lines.append(f"== 嵌入图（{total_embedded} 张，按来源文档分组） ==")
    for src, n in embedded_pairs:
        lines.append(f"  [{n}] 来自 {src}")
    lines.append("")
    lines.append("引用示例：evidence_refs.attachment = \".embedded_media/<原文档>::<图名>.<ext>\"")
    return "\n".join(lines)


def _refresh_embedded_media_index(
    *,
    root: Path,
    source_paths: List[Path],
    embedded_root: Path,
    items: List[AttachmentFile],
    by_filename: Dict[str, List[AttachmentFile]],
    by_rel_path: Dict[str, List[AttachmentFile]],
    status_counts: Dict[str, int],
    source_rel_path_by_logical_path: Dict[str, str],
) -> None:
    """Extract embedded images from office documents and mirror them as virtual items.

    Each source document is processed independently: a failure in one file never
    aborts extraction for the others. On-disk files are written before the
    virtual AttachmentFile is registered, and stale extracted files (belonging to
    source documents that no longer exist or were not extracted this run) are
    removed so the snapshot stays consistent with the index.
    """
    embedded_root.mkdir(parents=True, exist_ok=True)

    # Track which virtual rel_paths should exist after this run, and map disk
    # paths back to virtual rels so stale-file cleanup stays exact for nested
    # documents that share the same filename.
    expected_rels: set[str] = set()
    disk_to_rel: dict[str, str] = {}
    extracted_count = 0
    failed_sources: List[str] = []

    for path in source_paths:
        suffix = path.suffix.lower()
        if suffix == ".docx":
            extractor = extract_docx_media
        elif suffix == ".pptx":
            extractor = extract_pptx_media
        elif suffix == ".pdf":
            extractor = extract_pdf_media
        else:
            continue

        source_rel_path = path.relative_to(root)
        rel_source = source_rel_path.as_posix()
        try:
            media_items = extractor(path)
        except Exception as exc:
            _logger.error(
                "embedded media extraction failed for %s: %s", rel_source, exc,
                exc_info=_logger.isEnabledFor(logging.DEBUG),
            )
            failed_sources.append(rel_source)
            continue

        for m in media_items:
            disk_name = f"{path.name}__{m.media_filename}"
            disk_relative = source_rel_path.parent / disk_name
            out_path = embedded_root / disk_relative
            rel = f".embedded_media/{rel_source}::{m.media_filename}"
            try:
                out_path.parent.mkdir(parents=True, exist_ok=True)
                out_path.write_bytes(m.bytes)
            except Exception as exc:
                _logger.error(
                    "failed to write embedded media %s: %s", out_path, exc,
                    exc_info=_logger.isEnabledFor(logging.DEBUG),
                )
                failed_sources.append(f"{rel_source}/{m.media_filename}")
                continue

            expected_rels.add(rel)
            disk_to_rel[disk_relative.as_posix()] = rel
            source_rel_path_by_logical_path[_normalize_rel_path(rel)] = (
                (Path(".embedded_media") / disk_relative).as_posix()
            )

            existing = [it for it in items if it.rel_path == rel]
            if existing:
                continue

            item = AttachmentFile(
                index="",
                rel_dir=".embedded_media",
                filename=f"{path.name}::{m.media_filename}",
                rel_path=rel,
                file_type=m.file_type,
                description="",
                status="binary",
                extraction_status="binary",
                extracted_text="",
                size=len(m.bytes),
            )
            items.append(item)
            by_filename[item.filename.lower()].append(item)
            by_rel_path[_normalize_rel_path(item.rel_path)].append(item)
            status_counts["binary"] += 1
            extracted_count += 1

    # Remove on-disk extracted files that no longer match any registered virtual
    # item. This keeps the snapshot consistent when source documents change or
    # when a prior extraction run left stale files behind.
    if embedded_root.is_dir():
        for child in list(embedded_root.rglob("*")):
            if not child.is_file():
                continue
            disk_relative = child.relative_to(embedded_root).as_posix()
            virtual_rel = disk_to_rel.get(disk_relative)
            if virtual_rel is None or virtual_rel not in expected_rels:
                try:
                    child.unlink()
                    _logger.debug("removed stale embedded media file: %s", child)
                except Exception as exc:
                    _logger.warning("could not remove stale embedded media %s: %s", child, exc)

    # Drop virtual items whose on-disk file is no longer present.
    stale_items = [it for it in items if it.rel_path.startswith(".embedded_media/") and it.rel_path not in expected_rels]
    if stale_items:
        stale_rel_paths = {it.rel_path for it in stale_items}
        stale_filenames = {it.filename for it in stale_items}
        items[:] = [it for it in items if it.rel_path not in stale_rel_paths]
        for filename in stale_filenames:
            key = filename.lower()
            by_filename[key] = [it for it in by_filename.get(key, []) if it.rel_path not in stale_rel_paths]
            if not by_filename[key]:
                del by_filename[key]
        for rel_path in stale_rel_paths:
            key = _normalize_rel_path(rel_path)
            if key in by_rel_path:
                del by_rel_path[key]
            source_rel_path_by_logical_path.pop(key, None)
        status_counts["binary"] = max(0, status_counts.get("binary", 0) - len(stale_items))

    if extracted_count or failed_sources:
        _logger.info(
            "embedded media index refreshed: root=%s extracted=%d failed=%d",
            root, extracted_count, len(failed_sources),
        )
    if failed_sources:
        _logger.warning("embedded media extraction failures: %s", failed_sources)


def build_attachment_index(attachments_dir: str) -> Dict[str, object]:
    """Recursively index an attachment directory and extract bounded text.

    The returned maps are intentionally plain dictionaries so the index can be
    passed through the existing review pipeline without introducing a storage
    dependency. Unsupported regular files are retained as metadata with an
    ``unsupported`` extraction status so the constrained Agent can report them
    without guessing their contents.
    """
    if not attachments_dir:
        return {}
    root = Path(attachments_dir).expanduser()
    if not root.is_dir():
        return {}

    paths = [
        path for path in root.rglob("*")
        if path.is_file()
        and not path.is_symlink()
        and ".embedded_media" not in path.relative_to(root).parts
        and path.name not in _IGNORED_ATTACHMENT_NAMES
        and not path.name.startswith("~$")
        and ".embedded_media/" not in path.relative_to(root).as_posix()
        and path.relative_to(root).parts[:1] != (".embedded_media",)
    ]
    paths.sort(key=lambda path: path.relative_to(root).as_posix().lower())
    paths = paths[:_MAX_ATTACHMENT_FILES]

    items: List[AttachmentFile] = []
    by_filename: Dict[str, List[AttachmentFile]] = defaultdict(list)
    by_rel_path: Dict[str, List[AttachmentFile]] = defaultdict(list)
    by_index: Dict[str, List[AttachmentFile]] = defaultdict(list)
    by_sheet_norm: Dict[str, List[AttachmentFile]] = defaultdict(list)
    status_counts: Dict[str, int] = defaultdict(int)
    # `rel_path` is the stable path exposed to prompts, findings, and the UI.
    # Embedded entries intentionally use `::` there, while their on-disk names
    # use `__` so they remain portable to filesystems that reject colons.
    source_rel_path_by_logical_path: Dict[str, str] = {}

    for path in paths:
        relative_path = path.relative_to(root).as_posix()
        extracted_text, extraction_status = _extract_attachment_text(path)
        item = AttachmentFile(
            index=(_extract_indices(relative_path) or [""])[0],
            rel_dir=path.relative_to(root).parent.as_posix() if path.parent != root else "",
            filename=path.name,
            rel_path=relative_path,
            file_type=path.suffix.lower().lstrip("."),
            description=_limit_text(extracted_text, 400),
            status=extraction_status,
            size=path.stat().st_size,
            extracted_text=extracted_text,
            extraction_status=extraction_status,
        )
        items.append(item)
        by_filename[item.filename.lower()].append(item)
        normalized_path = _normalize_rel_path(item.rel_path)
        source_rel_path_by_logical_path[normalized_path] = item.rel_path
        path_parts = normalized_path.split("/")
        for start in range(len(path_parts)):
            suffix = "/".join(path_parts[start:])
            if item not in by_rel_path[suffix]:
                by_rel_path[suffix].append(item)
        for index in _extract_indices(relative_path):
            by_index[index].append(item)
        for sheet_norm in _extract_sheet_norms(item.rel_dir, item.rel_path):
            by_sheet_norm[sheet_norm].append(item)
        status_counts[extraction_status] += 1

    # After real attachment scanning, extract embedded media from DOCX/PPTX/PDF
    # and add them as virtual attachments indexed for Agent OCR.
    embedded_root = root / ".embedded_media"
    _refresh_embedded_media_index(
        root=root,
        source_paths=paths,
        embedded_root=embedded_root,
        items=items,
        by_filename=by_filename,
        by_rel_path=by_rel_path,
        status_counts=status_counts,
        source_rel_path_by_logical_path=source_rel_path_by_logical_path,
    )

    return {
        "path": str(root),
        "source_type": "directory",
        "items": items,
        "by_filename": dict(by_filename),
        "by_rel_path": dict(by_rel_path),
        "by_index": dict(by_index),
        "by_sheet_norm": dict(by_sheet_norm),
        "status_counts": dict(status_counts),
        "source_rel_path_by_logical_path": source_rel_path_by_logical_path,
    }


def _extract_attachment_refs(text: str) -> Tuple[List[str], List[str], List[str]]:
    s = (text or "").strip()
    if not s:
        return [], [], []
    rel_paths = [m.group(1) for m in ATTACHMENT_PATH_RE.finditer(s)]
    filenames = [m.group(1) for m in ATTACHMENT_FILE_RE.finditer(s)]
    indices = [m.group(1) for m in ATTACHMENT_INDEX_RE.finditer(s)]
    return filenames, rel_paths, indices


def _extract_cross_sheet_refs(text: str) -> List[str]:
    """Return the sheet identifiers mentioned by `<C22.XX-...>` or `C22.XX-...` refs."""
    s = (text or "").strip()
    if not s:
        return []
    seen: List[str] = []
    seen_norm: set = set()
    def _consume(ident: str) -> None:
        ident = (ident or "").strip().replace(" ", "")
        if not ident:
            return
        norm = _normalize_sheet_id(ident)
        if not norm or norm in seen_norm:
            return
        seen.append(ident)
        seen_norm.add(norm)
    for m in CROSS_SHEET_REF_RE.finditer(s):
        _consume(str(m.group(1) or ""))
    for m in CROSS_SHEET_BARE_RE.finditer(s):
        _consume(str(m.group(1) or ""))
    return seen


def _resolve_cross_sheet_items(
    attachments: Dict[str, object],
    sheet_refs: Sequence[str],
    *,
    limit: int = 6,
    workpaper_context: str = "",
) -> List[AttachmentFile]:
    """Look up attachment items that live under another sheet's directory.

    Used when the current sheet text references a sibling sheet's evidence
    (e.g. PE-6 cites `<C22.SA-3-1> 备份制度` but the file lives under
    ``审计证据/SA-3-1/SOP-65036-0 备份制度.pdf``). Without this, the LLM
    context for PE-6 sees only PE-6 attachments and misses the cross-sheet
    policy that design-effectiveness judgements actually depend on.

    If ``workpaper_context`` contains document titles (e.g. ``《备份制度》``),
    the resolver prefers items whose filename contains that title so the most
    relevant file bubbles to the top instead of being diluted by unrelated
    sibling attachments (a SA-3-1 directory typically holds dozens of docs).
    """
    if not attachments or not sheet_refs:
        return []
    by_sheet_norm = attachments.get("by_sheet_norm") or {}
    if not isinstance(by_sheet_norm, dict):
        return []
    # Extract title keywords (《...》) from the workpaper text. These are the
    # titles the auditor actually meant when writing `<C22.SA-3-1>`.
    title_keywords = re.findall(r"《([^》\s]{2,24})》", workpaper_context or "")
    picked: List[AttachmentFile] = []
    picked_keys: set = set()
    for ref in sheet_refs:
        norm = _normalize_sheet_id(ref)
        if not norm:
            continue
        items = by_sheet_norm.get(norm)
        if not isinstance(items, list) or not items:
            continue
        # Sort by relevance: title keyword hits first, otherwise basename order.
        def _score(it: AttachmentFile) -> Tuple[int, str]:
            name = (it.rel_path or it.filename or "").lower()
            hits = sum(1 for kw in title_keywords if kw and kw.lower() in name)
            return (-hits, name)
        ordered = sorted(items, key=_score) if items else items
        for it in ordered:
            if not isinstance(it, AttachmentFile):
                continue
            key = (it.rel_path or it.filename or "").lower()
            if not key or key in picked_keys:
                continue
            picked_keys.add(key)
            picked.append(it)
            if len(picked) >= limit:
                return picked
    return picked


def _compact_keywords(text: str) -> List[str]:
    s = (text or "").strip()
    if not s:
        return []
    tokens = re.findall(r"[一-鿿]{2,}|[A-Za-z0-9]{2,}", s)
    stop = {"审计", "程序", "执行", "标准", "附件", "证据", "截图", "导出", "清单", "日志", "台账"}
    out: List[str] = []
    seen = set()
    for t in tokens:
        tt = t.strip().lower()
        if not tt or tt in stop:
            continue
        if tt in seen:
            continue
        seen.add(tt)
        out.append(t.strip())
        if len(out) >= 18:
            break
    return out


def _evidence_matches_step(step_text: str, attachment_desc: str) -> bool:
    if not step_text or not attachment_desc:
        return True
    step_keys = set(_compact_keywords(step_text) + [k for k in CHECKPOINT_VOCAB if k in step_text])
    desc_keys = set(_compact_keywords(attachment_desc) + [k for k in CHECKPOINT_VOCAB if k in attachment_desc])
    if not step_keys or not desc_keys:
        return True
    inter = step_keys.intersection(desc_keys)
    return len(inter) >= 1


def _extract_sheet_norms(*texts: str) -> List[str]:
    joined = " ".join(str(t or "") for t in texts if t)
    if not joined:
        return []
    norms: List[str] = []
    for m in _SHEET_TAG_RE.findall(joined) + _SHEET_TAG_NODELIM_RE.findall(joined):
        nm = _normalize_sheet_id(m)
        if nm and nm not in norms:
            norms.append(nm)
    return norms


def load_attachments_preview_xlsx(preview_path: str) -> Dict[str, object]:
    if not preview_path:
        return {}
    wb = openpyxl.load_workbook(preview_path, data_only=True)
    ws = wb["图片描述"] if "图片描述" in wb.sheetnames else wb.active

    header_row = 1
    header_map: Dict[str, int] = {}
    max_scan_col = min(ws.max_column or 0, 40)
    for c in range(1, max_scan_col + 1):
        v = ws.cell(row=header_row, column=c).value
        if not v:
            continue
        key = str(v).strip().replace("\n", "").replace(" ", "")
        if key:
            header_map[key] = c

    def _col(*names: str) -> int:
        for n in names:
            n2 = str(n).strip().replace("\n", "").replace(" ", "")
            if n2 in header_map:
                return int(header_map[n2])
        return 0

    col_index = _col("目录索引", "索引")
    col_rel_dir = _col("相对目录", "目录", "文件夹", "相对文件夹")
    col_filename = _col("附件文件名", "文件名", "附件名称", "图片文件名")
    col_rel_path = _col("相对路径", "路径", "附件路径", "图片路径")
    col_file_type = _col("文件类型", "类型", "后缀", "格式")
    col_desc = _col("详细描述", "描述", "内容", "图片描述")
    col_status = _col("状态", "校验状态", "结果")

    items: List[AttachmentPreviewItem] = []
    by_filename: Dict[str, List[AttachmentPreviewItem]] = defaultdict(list)
    by_rel_path: Dict[str, List[AttachmentPreviewItem]] = defaultdict(list)
    by_index: Dict[str, List[AttachmentPreviewItem]] = defaultdict(list)
    by_sheet_norm: Dict[str, List[AttachmentPreviewItem]] = defaultdict(list)
    status_counts: Dict[str, int] = defaultdict(int)

    for r in range(2, (ws.max_row or 0) + 1):
        raw_filename = ws.cell(row=r, column=col_filename).value if col_filename else None
        raw_desc = ws.cell(row=r, column=col_desc).value if col_desc else None
        raw_rel_path = ws.cell(row=r, column=col_rel_path).value if col_rel_path else None
        raw_status = ws.cell(row=r, column=col_status).value if col_status else None
        raw_index = ws.cell(row=r, column=col_index).value if col_index else None
        raw_rel_dir = ws.cell(row=r, column=col_rel_dir).value if col_rel_dir else None
        raw_file_type = ws.cell(row=r, column=col_file_type).value if col_file_type else None

        filename = str(raw_filename).strip() if raw_filename else ""
        description = str(raw_desc).strip() if raw_desc else ""
        rel_path = str(raw_rel_path).strip() if raw_rel_path else ""
        status = str(raw_status).strip() if raw_status else ""
        index = str(raw_index).strip() if raw_index is not None else ""
        rel_dir = str(raw_rel_dir).strip() if raw_rel_dir else ""
        file_type = str(raw_file_type).strip() if raw_file_type else ""

        if not filename and not rel_path and not description:
            continue

        item = AttachmentPreviewItem(
            index=index,
            rel_dir=rel_dir,
            filename=filename,
            rel_path=rel_path,
            file_type=file_type,
            description=description,
            status=status,
        )
        items.append(item)

        if filename:
            by_filename[filename.lower()].append(item)
        if rel_path:
            by_rel_path[rel_path.lower().replace("/", "\\")].append(item)
        if index:
            by_index[index].append(item)
        for sn in _extract_sheet_norms(rel_dir, rel_path):
            by_sheet_norm[sn].append(item)
        status_counts[status or ""] += 1

    return {
        "path": preview_path,
        "source_type": "preview",
        "items": items,
        "by_filename": dict(by_filename),
        "by_rel_path": dict(by_rel_path),
        "by_index": dict(by_index),
        "by_sheet_norm": dict(by_sheet_norm),
        "status_counts": dict(status_counts),
    }


def _match_attachment_items(
    attachments: Dict[str, object],
    *,
    filenames: Sequence[str],
    rel_paths: Sequence[str],
    indices: Sequence[str],
) -> Tuple[List[AttachmentFile], List[str]]:
    if not attachments:
        return [], list(indices) + list(rel_paths) + list(filenames)
    by_filename = attachments.get("by_filename") or {}
    by_rel_path = attachments.get("by_rel_path") or {}
    by_index = attachments.get("by_index") or {}

    picked: List[AttachmentFile] = []
    picked_keys = set()
    missing: List[str] = []

    def _add(items: Iterable[AttachmentFile]) -> None:
        for it in items:
            key = (it.rel_path or "").lower() or (it.filename or "").lower()
            if not key:
                continue
            if key in picked_keys:
                continue
            picked_keys.add(key)
            picked.append(it)

    for idx in indices:
        lst = by_index.get(str(idx).strip())
        if isinstance(lst, list) and lst:
            _add(lst)
        else:
            missing.append(f"索引{idx}")

    for p in rel_paths:
        key = _normalize_rel_path(str(p))
        lst = by_rel_path.get(key)
        if not lst:
            lst = by_rel_path.get(key.replace("/", "\\"))
        if isinstance(lst, list) and lst:
            _add(lst)
            continue
        missing.append(p)

    for f in filenames:
        key = str(f).strip().lower()
        lst = by_filename.get(key)
        if isinstance(lst, list) and lst:
            _add(lst)
            continue
        missing.append(f)

    return picked, missing


def _match_preview_items(
    preview: Dict[str, object],
    *,
    filenames: Sequence[str],
    rel_paths: Sequence[str],
    indices: Sequence[str],
) -> Tuple[List[AttachmentFile], List[str]]:
    """Backward-compatible alias for older callers and preview fixtures."""
    return _match_attachment_items(
        preview,
        filenames=filenames,
        rel_paths=rel_paths,
        indices=indices,
    )


def _verified_attachment_text(attachments: Dict[str, object], item: AttachmentFile) -> str:
    """Return deterministic or cached OCR text for one indexed attachment."""
    extracted = str(getattr(item, "extracted_text", "") or "")
    if extracted:
        return extracted
    key = _normalize_rel_path(str(item.rel_path or item.filename or ""))
    ocr_cache = attachments.get("ocr_by_path") or {}
    cached = ocr_cache.get(key) if isinstance(ocr_cache, dict) else None
    if isinstance(cached, dict) and str(cached.get("status", "")).lower() == "ok":
        return str(cached.get("content", "") or "")
    return ""


def _verify_attachment_evidence_refs(
    evidence_refs: Sequence[dict],
    attachments: Dict[str, object],
    ws=None,
) -> List[dict]:
    """Keep attachment references only when their source path and excerpt verify.

    Cell evidence is retained when an LLM added an invalid attachment hint, but
    the unverified attachment field is removed. This prevents a hallucinated
    file path from surviving into a structured Finding while preserving a valid
    workbook-cell citation.
    """
    from review.excel_utils import _get_cell_value
    from review.validation import _verify_evidence_refs

    verified: List[dict] = []
    if not attachments:
        # An attachment citation cannot be verified without the pinned index.
        # Preserve only an independently verifiable workbook-cell citation;
        # otherwise fail closed rather than exposing an LLM-supplied filename.
        for raw in evidence_refs:
            if not isinstance(raw, dict):
                continue
            fallback = dict(raw)
            had_attachment = bool(str(fallback.pop("attachment", "") or "").strip())
            if not had_attachment:
                verified.append(fallback)
                continue
            if not fallback.get("cell_or_range"):
                continue
            if ws is None:
                verified.append(fallback)
            else:
                verified.extend(_verify_evidence_refs([fallback], ws))
        return verified

    for raw in evidence_refs:
        if not isinstance(raw, dict):
            continue
        attachment = str(raw.get("attachment", "") or "").strip()
        if not attachment:
            verified.append(dict(raw))
            continue
        normalized_attachment = _normalize_rel_path(attachment)
        indexed_paths = attachments.get("by_rel_path") or {}
        if isinstance(indexed_paths, dict) and indexed_paths.get(normalized_attachment):
            # Virtual embedded-media paths contain `::`; the generic filename
            # regex would otherwise reduce them to their trailing image name.
            filenames, rel_paths, indices = [], [attachment], []
        else:
            filenames, rel_paths, indices = _extract_attachment_refs(attachment)
            if not filenames and not rel_paths and not indices:
                rel_paths = [attachment]
        matched, _ = _match_attachment_items(
            attachments,
            filenames=filenames,
            rel_paths=rel_paths,
            indices=indices,
        )
        excerpt = str(raw.get("excerpt", "") or "").strip()
        excerpt_matches = [
            item
            for item in matched
            if excerpt and excerpt in _verified_attachment_text(attachments, item)
        ]
        # An attachment name alone does not prove what was inspected. Keep it
        # only when the finding carries a non-empty quote found in the pinned
        # source text (including cached OCR text).
        # A basename can resolve to multiple files. Pin the citation only when
        # exactly one indexed source contains the quoted text; otherwise its
        # source is ambiguous and must not survive as attachment evidence.
        if not excerpt or len(excerpt_matches) != 1:
            fallback = dict(raw)
            fallback.pop("attachment", None)
            if fallback.get("cell_or_range"):
                if ws is None:
                    verified.append(fallback)
                else:
                    verified.extend(_verify_evidence_refs([fallback], ws))
            continue
        cell = str(raw.get("cell_or_range", "") or "").strip()
        if ws is not None and cell and not _get_cell_value(ws, cell):
            continue
        normalized = dict(raw)
        matched_item = excerpt_matches[0]
        normalized["attachment"] = str(
            matched_item.rel_path or matched_item.filename or attachment
        )
        verified.append(normalized)
    return verified


def format_evidence_refs_for_basis(
    evidence_refs: Sequence[dict],
    *,
    max_refs: int = 3,
    max_excerpt_chars: int = 200,
) -> str:
    """Render verified evidence as readable citations in a finding basis."""
    rendered: List[str] = []
    for ref in evidence_refs[:max(0, max_refs)]:
        if not isinstance(ref, dict):
            continue
        excerpt = str(ref.get("excerpt", "") or "").strip()
        if not excerpt:
            continue
        attachment = str(ref.get("attachment", "") or "").strip()
        if attachment:
            location = f"附件：{attachment}"
        else:
            location = "!".join(
                value for value in (
                    str(ref.get("sheet", "") or "").strip(),
                    str(ref.get("cell_or_range", "") or "").strip(),
                ) if value
            )
        rendered.append(f"{location or '未定位'}: {excerpt[:max_excerpt_chars]}")
    return "; ".join(rendered)


def _attachments_context_for_sheet(ws, attachments: Dict[str, object], limit_chars: int = 12000) -> str:
    if not attachments:
        return ""
    text = _build_sheet_text_for_llm(ws, max_cells=260, max_chars=24000)
    filenames, rel_paths, indices = _extract_attachment_refs(text)
    matched, _ = _match_attachment_items(attachments, filenames=filenames, rel_paths=rel_paths, indices=indices)
    # Cross-sheet evidence (e.g. PE-6 citing `<C22.SA-3-1> 备份制度`). We pull
    # the first few items from the referenced sibling sheet directory so the
    # LLM can verify design-effectiveness claims against the actual policy
    # document, not just the C8 cell narrative that paraphrases it. The
    # resolver prefers files whose name matches a `《...》` title mentioned in
    # the same workpaper (so `《备份制度》` bubbles the SOP-65036 PDF to the
    # top of the SA-3-1 directory).
    sheet_refs = _extract_cross_sheet_refs(text)
    for it in _resolve_cross_sheet_items(attachments, sheet_refs, limit=4, workpaper_context=text):
        matched.append(it)
    by_sheet = attachments.get("by_sheet_norm") or {}
    if isinstance(by_sheet, dict):
        norm = _normalize_sheet_id(getattr(ws, "title", "") or "")
        lst = by_sheet.get(norm)
        if isinstance(lst, list) and lst:
            seen = set(id(it) for it in matched)
            for it in lst[:80]:
                if not isinstance(it, AttachmentPreviewItem):
                    continue
                if id(it) in seen:
                    continue
                matched.append(it)
                seen.add(id(it))
                if len(matched) >= 80:
                    break
    agent_by_sheet = attachments.get("agent_evidence_by_sheet") or {}
    agent_evidence = []
    if isinstance(agent_by_sheet, dict):
        candidate = agent_by_sheet.get(_normalize_sheet_id(getattr(ws, "title", "") or ""))
        if isinstance(candidate, list):
            agent_evidence = [item for item in candidate if isinstance(item, dict)]
    if not matched and not agent_evidence:
        return ""
    parts: List[str] = []
    total = 0
    for evidence in agent_evidence[:20]:
        path = str(evidence.get("path", "") or "")
        excerpt = str(evidence.get("excerpt", "") or "")
        if not path or not excerpt:
            continue
        line = (
            f"- Agent调查证据 | {path} | 类型={evidence.get('file_type', '')} | "
            f"解析状态={evidence.get('extraction_status', 'unknown')}"
            f"\n  原文摘录：{_truncate(excerpt, 1800)}"
        )
        supports = str(evidence.get("supports", "") or "").strip()
        if supports:
            line += f"\n  调查用途：{_truncate(supports, 300)}"
        if total + len(line) + 1 > limit_chars:
            break
        parts.append(line)
        total += len(line) + 1
    for it in matched[:40]:
        verified_text = _verified_attachment_text(attachments, it)
        display_status = it.extraction_status or it.status or "unknown"
        if verified_text and str(display_status).lower() not in {"ok", "ocr"}:
            display_status = "ocr"
        line = (
            f"- {it.rel_path or it.filename} | 类型={it.file_type} | 大小={it.size} | "
            f"解析状态={display_status}"
        )
        if verified_text:
            line += f"\n  内容：{_truncate(verified_text, 1800)}"
        elif it.description:
            line += f"\n  描述：{it.description}"
        if total + len(line) + 1 > limit_chars:
            break
        parts.append(line)
        total += len(line) + 1
    return "\n".join(parts).strip()


def _check_attachment_references(ws_title: str, ws, attachments: Dict[str, object]) -> List[Finding]:
    if not attachments:
        return []
    findings: List[Finding] = []

    used_any = False
    for coord, text in _extract_sheet_text_cells(ws):
        filenames, rel_paths, indices = _extract_attachment_refs(text)
        if not filenames and not rel_paths and not indices:
            continue
        used_any = True
        matched, missing = _match_attachment_items(
            attachments,
            filenames=filenames,
            rel_paths=rel_paths,
            indices=indices,
        )
        if missing:
            findings.append(
                Finding(
                    issue_type="附件证据引用未匹配到附件目录",
                    severity="P2",
                    sheet=ws_title,
                    cell=coord,
                    snippet=_truncate(text, 220),
                    basis=_truncate(
                        "引用: "
                        + "、".join(sorted({m for m in missing if m}))
                        + f"\n附件目录: {attachments.get('path','')}",
                        1200,
                    ),
                    suggestion="核对底稿中的附件编号/文件名/路径是否与附件目录一致；如存在重命名或遗漏，请补齐目录或更新引用。",
                    **deterministic_finding_fields(
                        origin="attachment_reference",
                        rule_hint="attachment_reference_missing",
                        assertion_id="attachment.inventory.presence",
                        sheet=ws_title,
                        cell=coord,
                        claim_subject=(
                            f"{ws_title}|attachment:{sorted({m for m in missing if m})[0]}"
                            if any(m for m in missing)
                            else ""
                        ),
                        claim_value="absent",
                    ),
                )
            )
        bad = [
            it for it in matched
            if (it.extraction_status or it.status) not in {"", "ok", "OK"}
            and not _verified_attachment_text(attachments, it)
        ]
        for it in bad[:5]:
            findings.append(
                Finding(
                    issue_type="附件证据内容未解析",
                    severity="P1",
                    sheet=ws_title,
                    cell=coord,
                    snippet=_truncate(text, 220),
                    basis=_truncate(
                        f"附件: {it.rel_path or it.filename}\n"
                        f"解析状态: {it.extraction_status or it.status}\n"
                        f"类型: {it.file_type}",
                        1200,
                    ),
                    suggestion="检查该附件是否为图片、扫描件或不支持的格式；必要时补充可检索版本，或配置 OCR/视觉模型后重新审阅。",
                    **deterministic_finding_fields(
                        origin="attachment_reference",
                        rule_hint="attachment_text_unavailable",
                        assertion_id="attachment.content.support",
                        sheet=ws_title,
                        cell=coord,
                        claim_subject=f"{ws_title}|attachment:{it.rel_path or it.filename}",
                        claim_value="unavailable",
                    ),
                )
            )

    if not used_any:
        return []
    return findings
